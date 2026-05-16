from __future__ import annotations

import hashlib
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from document_inteligence.domain.entities import (
    DocumentElement,
    DocumentPage,
    ElementType,
    ParsedDocument,
)
from document_inteligence.domain.repositories import DocumentIngestor
from document_inteligence.domain.value_objects import BBox
from document_inteligence.infrastructure.ingestors.pptx_table_xml import (
    NS as _A_NS,
    ParsedTable,
    extract_tables_from_pptx,
    paragraphs_with_bullets,
    paragraphs_with_formatting,
)


MATH_NS = {"m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}

_EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 DPI


def _clamp_01(value: float) -> float:
    return max(0.0, min(1.0, value))


def _shape_type(shape) -> ElementType:
    if getattr(shape, "has_table", False):
        return ElementType.TABLE
    if getattr(shape, "has_chart", False):
        return ElementType.CHART
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ElementType.IMAGE
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return ElementType.GROUP
    if getattr(shape, "has_text_frame", False):
        return ElementType.TEXT
    return ElementType.UNKNOWN


def _extract_text(shape) -> str | None:
    if getattr(shape, "has_text_frame", False):
        raw_text = _extract_text_with_bullets(shape)
        math_exprs = _extract_math_expressions_from_shape(shape)
        math_lines = [f"[MATH] {expr}" for expr in math_exprs if expr]
        if raw_text and math_lines:
            return raw_text + "\n" + "\n".join(math_lines)
        if raw_text:
            return raw_text
        if math_lines:
            return "\n".join(math_lines)
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            values = [cell.text.strip() for cell in row.cells]
            line = " | ".join([v for v in values if v])
            if line:
                rows.append(line)
        raw = "\n".join(rows)
        return raw or None
    return None


def _extract_text_with_bullets(shape) -> str:
    """Extract text preserving bullet/numbering prefixes from paragraph XML."""
    element = getattr(shape, "_element", None)
    if element is None:
        return (shape.text or "").strip()
    p_elements = element.findall(".//{%s}p" % _A_NS["a"])
    if not p_elements:
        return (shape.text or "").strip()
    return paragraphs_with_bullets(p_elements, _A_NS).strip()


def _extract_formatted_text(shape) -> str | None:
    """Extract text with bold/italic/underline Markdown markers per run."""
    element = getattr(shape, "_element", None)
    if element is None:
        return None
    p_elements = element.findall(".//{%s}p" % _A_NS["a"])
    if not p_elements:
        return None
    result = paragraphs_with_formatting(p_elements, _A_NS, mode="markdown").strip()
    return result or None


def _build_bbox(shape, slide_width: float, slide_height: float) -> BBox:
    x = _clamp_01(float(shape.left) / slide_width)
    y = _clamp_01(float(shape.top) / slide_height)
    width = max(_clamp_01(float(shape.width) / slide_width), 1e-6)
    height = max(_clamp_01(float(shape.height) / slide_height), 1e-6)
    return BBox(x=x, y=y, width=width, height=height)


def _build_bbox_from_abs(
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    slide_width: float,
    slide_height: float,
) -> BBox:
    x = _clamp_01(float(left) / slide_width)
    y = _clamp_01(float(top) / slide_height)
    width = max(_clamp_01(float(width) / slide_width), 1e-6)
    height = max(_clamp_01(float(height) / slide_height), 1e-6)
    return BBox(x=x, y=y, width=width, height=height)


def _extract_math_expressions_from_shape(shape) -> list[str]:
    """Extract linearized math expressions from OMML nodes in a shape."""
    element = getattr(shape, "_element", None)
    if element is None:
        return []
    try:
        omaths = element.xpath(".//m:oMath | .//m:oMathPara", namespaces=MATH_NS)
    except Exception:
        return []

    expressions: list[str] = []
    for node in omaths:
        try:
            tokens = node.xpath(".//m:t/text()", namespaces=MATH_NS)
        except Exception:
            tokens = []
        expr = " ".join([t.strip() for t in tokens if str(t).strip()]).strip()
        if expr:
            expressions.append(expr)
    return expressions


class PptxIngestor(DocumentIngestor):
    def __init__(
        self,
        asset_output_dir: str | None = None,
        include_master_shapes: bool = True,
        deduplicate_master_shapes: bool = True,
    ):
        self._asset_output_dir = Path(asset_output_dir) if asset_output_dir else None
        self._include_master = include_master_shapes
        self._dedup_master = deduplicate_master_shapes

    def supports(self, path: str) -> bool:
        return Path(path).suffix.lower() == ".pptx"

    def ingest(self, path: str) -> ParsedDocument:
        presentation = Presentation(path)
        xml_tables_by_slide = extract_tables_from_pptx(Path(path))
        slide_width = float(presentation.slide_width)
        slide_height = float(presentation.slide_height)

        pages: list[DocumentPage] = []
        seen_master_hashes: set[str] = set()
        for slide_index, slide in enumerate(presentation.slides, start=1):
            elements: list[DocumentElement] = []
            slide_table_idx = 0
            slide_tables = xml_tables_by_slide.get(slide_index, [])
            element_counter = 0

            def process_shape(
                shape,
                abs_left: float | None = None,
                abs_top: float | None = None,
                source: str = "slide",
            ) -> None:
                nonlocal slide_table_idx, element_counter

                shape_type = _shape_type(shape)
                # group child coordinates are relative to parent group.
                left = float(shape.left) if abs_left is None else abs_left
                top = float(shape.top) if abs_top is None else abs_top
                width = float(shape.width)
                height = float(shape.height)

                metadata: dict[str, object] = {
                    "shape_name": getattr(shape, "name", None),
                    "shape_type": str(getattr(shape, "shape_type", "")),
                    "has_table": bool(getattr(shape, "has_table", False)),
                    "has_chart": bool(getattr(shape, "has_chart", False)),
                    "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                    "source": source,
                }

                if shape_type == ElementType.TABLE:
                    parsed_table: ParsedTable | None = None
                    if source == "slide" and slide_table_idx < len(slide_tables):
                        parsed_table = slide_tables[slide_table_idx]
                    if source == "slide":
                        slide_table_idx += 1
                    if parsed_table is not None:
                        cell_bboxes = _build_table_cell_bboxes(
                            parsed_table=parsed_table,
                            left=left,
                            top=top,
                            width=width,
                            height=height,
                            slide_width=slide_width,
                            slide_height=slide_height,
                        )
                        metadata["table_col_widths"] = parsed_table.col_widths
                        metadata["table_row_heights"] = parsed_table.row_heights
                        metadata["table_cells"] = [
                            {
                                "row": cell.row,
                                "col": cell.col,
                                "text": cell.text,
                                "is_spanned": cell.is_spanned,
                                "span_width": cell.span_width,
                                "span_height": cell.span_height,
                                "is_merge_origin": cell.is_merge_origin,
                                "cell_bbox": cell_bboxes.get((cell.row, cell.col)),
                                "formatted_text": cell.formatted_text,
                            }
                            for cell in parsed_table.cells
                        ]

                element_counter += 1
                z_order = element_counter

                if shape_type == ElementType.IMAGE:
                    image_info = self._extract_and_map_image(path=path, slide_index=slide_index, z_order=z_order, shape=shape)
                    if image_info:
                        metadata["image"] = image_info

                math_exprs = _extract_math_expressions_from_shape(shape)
                if math_exprs:
                    metadata["math_expressions"] = math_exprs

                if shape_type == ElementType.TEXT:
                    fmt_text = _extract_formatted_text(shape)
                    if fmt_text:
                        metadata["formatted_text"] = fmt_text

                element = DocumentElement(
                    element_id=f"E_{slide_index:03d}_{z_order:04d}",
                    element_type=shape_type,
                    page_number=slide_index,
                    z_order=z_order,
                    bbox=_build_bbox_from_abs(
                        left=left,
                        top=top,
                        width=width,
                        height=height,
                        slide_width=slide_width,
                        slide_height=slide_height,
                    ),
                    text=_extract_text(shape),
                    metadata=metadata,
                )
                elements.append(element)

                if shape_type == ElementType.GROUP:
                    for child in shape.shapes:
                        child_abs_left = left + float(child.left)
                        child_abs_top = top + float(child.top)
                        process_shape(child, abs_left=child_abs_left, abs_top=child_abs_top, source=source)

            for shape in slide.shapes:
                process_shape(shape, source="slide")

            if self._include_master:
                slide_ph_idxs = _slide_placeholder_idxs(slide)
                layout = slide.slide_layout
                _ingest_extra_shapes(
                    layout, "layout", slide_ph_idxs,
                    process_shape, self._dedup_master, seen_master_hashes,
                )
                try:
                    master = layout.slide_master
                except Exception:
                    master = None
                if master is not None:
                    _ingest_extra_shapes(
                        master, "master", slide_ph_idxs,
                        process_shape, self._dedup_master, seen_master_hashes,
                    )

            pages.append(
                DocumentPage(
                    page_number=slide_index,
                    width=slide_width,
                    height=slide_height,
                    elements=elements,
                )
            )

        return ParsedDocument(source_path=path, pages=pages)

    def _extract_and_map_image(self, path: str, slide_index: int, z_order: int, shape) -> dict[str, object] | None:
        try:
            image = shape.image
            ext = image.ext or "bin"
            blob = image.blob
        except Exception:
            return None

        blob = _apply_picture_crop_if_needed(shape=shape, blob=blob, ext=ext)

        display_w_px = round(float(shape.width) / _EMU_PER_PX)
        display_h_px = round(float(shape.height) / _EMU_PER_PX)

        filename = f"slide{slide_index:03d}_shape{z_order:04d}.{ext}"
        info: dict[str, object] = {
            "filename": filename,
            "ext": ext,
            "bytes": len(blob),
            "display_width_px": display_w_px,
            "display_height_px": display_h_px,
        }
        if self._asset_output_dir is None:
            return info

        self._asset_output_dir.mkdir(parents=True, exist_ok=True)
        target = self._asset_output_dir / filename
        target.write_bytes(blob)
        # Markdown is typically saved near outputs; keep asset link portable by
        # using "assets_dir_name/filename" instead of cwd-dependent full path.
        info["relative_path"] = f"{self._asset_output_dir.name}/{filename}"
        return info


def _apply_picture_crop_if_needed(*, shape, blob: bytes, ext: str) -> bytes:
    crop_left = float(getattr(shape, "crop_left", 0.0) or 0.0)
    crop_right = float(getattr(shape, "crop_right", 0.0) or 0.0)
    crop_top = float(getattr(shape, "crop_top", 0.0) or 0.0)
    crop_bottom = float(getattr(shape, "crop_bottom", 0.0) or 0.0)

    if crop_left == 0.0 and crop_right == 0.0 and crop_top == 0.0 and crop_bottom == 0.0:
        return blob

    try:
        with Image.open(BytesIO(blob)) as img:
            width, height = img.size
            left = int(max(0.0, min(1.0, crop_left)) * width)
            right_trim = int(max(0.0, min(1.0, crop_right)) * width)
            top = int(max(0.0, min(1.0, crop_top)) * height)
            bottom_trim = int(max(0.0, min(1.0, crop_bottom)) * height)

            right = max(left + 1, width - right_trim)
            bottom = max(top + 1, height - bottom_trim)
            if right <= left or bottom <= top:
                return blob

            cropped = img.crop((left, top, right, bottom))
            out = BytesIO()
            save_format = _image_save_format(ext)
            cropped.save(out, format=save_format)
            return out.getvalue()
    except Exception:
        return blob


def _image_save_format(ext: str) -> str:
    e = ext.lower()
    if e in {"jpg", "jpeg"}:
        return "JPEG"
    if e == "gif":
        return "GIF"
    if e == "bmp":
        return "BMP"
    if e == "tiff":
        return "TIFF"
    return "PNG" if e == "png" else "PNG"


def _build_table_cell_bboxes(
    *,
    parsed_table: ParsedTable,
    left: float,
    top: float,
    width: float,
    height: float,
    slide_width: float,
    slide_height: float,
) -> dict[tuple[int, int], dict[str, float]]:
    col_widths = parsed_table.col_widths or [1.0]
    row_heights = parsed_table.row_heights or [1.0]
    total_col = sum(col_widths) if sum(col_widths) > 0 else float(len(col_widths))
    total_row = sum(row_heights) if sum(row_heights) > 0 else float(len(row_heights))

    col_starts = [0.0]
    for w in col_widths:
        col_starts.append(col_starts[-1] + float(w) / total_col)
    row_starts = [0.0]
    for h in row_heights:
        row_starts.append(row_starts[-1] + float(h) / total_row)

    out: dict[tuple[int, int], dict[str, float]] = {}
    for cell in parsed_table.cells:
        r = int(cell.row)
        c = int(cell.col)
        col_end_idx = min(c + max(1, int(cell.span_width)), len(col_starts) - 1)
        row_end_idx = min(r + max(1, int(cell.span_height)), len(row_starts) - 1)
        if c >= len(col_starts) - 1 or r >= len(row_starts) - 1:
            continue

        cell_left = left + (width * col_starts[c])
        cell_top = top + (height * row_starts[r])
        cell_width = max(1.0, width * (col_starts[col_end_idx] - col_starts[c]))
        cell_height = max(1.0, height * (row_starts[row_end_idx] - row_starts[r]))
        bbox = _build_bbox_from_abs(
            left=cell_left,
            top=cell_top,
            width=cell_width,
            height=cell_height,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        out[(r, c)] = {
            "x": bbox.x,
            "y": bbox.y,
            "width": bbox.width,
            "height": bbox.height,
        }
    return out


def _slide_placeholder_idxs(slide) -> set[int]:
    out: set[int] = set()
    try:
        for ph in slide.placeholders:
            out.add(ph.placeholder_format.idx)
    except Exception:
        pass
    return out


def _ingest_extra_shapes(
    source_obj,
    source_label: str,
    slide_ph_idxs: set[int],
    process_shape_fn,
    dedup: bool,
    seen_hashes: set[str],
) -> None:
    try:
        shapes = source_obj.shapes
    except Exception:
        return
    for shape in shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                idx = shape.placeholder_format.idx
                if idx in slide_ph_idxs:
                    continue
            except Exception:
                pass

        text = ""
        try:
            text = (shape.text or "").strip()
        except Exception:
            pass
        if not text:
            if not getattr(shape, "has_table", False) and not _is_picture(shape):
                continue
        if _is_master_template_text(text):
            continue

        if dedup:
            content_hash = hashlib.md5(
                f"{source_label}:{getattr(shape, 'name', '')}:{text}".encode()
            ).hexdigest()
            if content_hash in seen_hashes:
                continue
            seen_hashes.add(content_hash)

        try:
            process_shape_fn(shape, source=source_label)
        except Exception:
            pass


def _is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


_MASTER_TEMPLATE_PATTERNS = {
    "click to edit master",
    "마스터 제목 스타일 편집",
    "마스터 텍스트 스타일 편집",
    "second level",
    "third level",
    "fourth level",
    "fifth level",
}


def _is_master_template_text(text: str) -> bool:
    if not text:
        return False
    low = text.lower().strip()
    if low in {"‹#›", "<#>", "‹날짜›", "<날짜>"}:
        return True
    for pattern in _MASTER_TEMPLATE_PATTERNS:
        if pattern in low:
            return True
    return False
